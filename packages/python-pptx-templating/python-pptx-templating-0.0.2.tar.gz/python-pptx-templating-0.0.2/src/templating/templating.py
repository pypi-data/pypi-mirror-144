import re

import chevron
from pptx import Presentation as PptxPresentation
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import PicturePlaceholder
from pptx.slide import Slide as PptxSlide, SlideLayout


def find_shape_by_name(slide: PptxSlide, shape_name: str):
    for shape in slide.shapes:
        if shape.name == shape_name:
            return shape


def format_key_value_shape(shape):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            key_value = key_value_pattern.search(r.text)
            if key_value:
                value_run = p.add_run()

                r.text = key_value.group("key")
                r.font.bold = True
                value_run.text = key_value.group("value")


key_value_pattern = re.compile(r"(?P<key>.*:)(?P<value>.*)", re.IGNORECASE)


def format_paragraph(paragraph):
    def strip_dashes(p):
        if p.runs:
            p.runs[0].text = p.runs[0].text.lstrip("- ")

    if paragraph.text.startswith("---"):
        paragraph.level = 3
        strip_dashes(paragraph)
    elif paragraph.text.startswith("--"):
        paragraph.level = 2
        strip_dashes(paragraph)
    elif paragraph.text.startswith("-"):
        paragraph.level = 1
        strip_dashes(paragraph)


def replace_shape_with_picture(slide: PptxSlide, shape: BaseShape, img):
    if isinstance(shape, Picture):
        replace_with_picture(slide, shape, img)
    elif isinstance(shape, PicturePlaceholder):
        replace_placeholder_with_image(slide, shape, img)
    else:
        raise ValueError(f"""Cannot replace {type(shape).__name__} with picture""")


def update_and_format_shape_text(shape, text):
    # TODO support a subset of MarkDown format
    # TODO support of "<b>Name:</b> value" formatting
    shape.text = text
    for p in shape.text_frame.paragraphs:
        format_paragraph(p)


def is_picture_variable(variable: str) -> bool:
    return ("picture" in variable) or ("image" in variable)


def replace_with_picture(slide: PptxSlide, picture: Picture, img):
    new_shape = slide.shapes.add_picture(
        img,
        picture.left,
        picture.top,
        picture.width,
        picture.height,
    )
    new_shape.auto_shape_type = picture.auto_shape_type
    old_pic = picture._element
    new_pic = new_shape._element
    old_pic.addnext(new_pic)
    old_pic.getparent().remove(old_pic)


def replace_placeholder_with_image(
    slide: PptxSlide, placeholder: PicturePlaceholder, img
):
    pic = slide.shapes.add_picture(img, placeholder.left, placeholder.top)

    # calculate max width/height for target size
    ratio = min(
        placeholder.width / float(pic.width), placeholder.height / float(pic.height)
    )

    pic.height = int(pic.height * ratio)
    pic.width = int(pic.width * ratio)

    pic.left = int(placeholder.left + ((placeholder.width - pic.width) / 2))
    pic.top = int(placeholder.top + ((placeholder.height - pic.height) / 2))

    p = placeholder.element
    p.getparent().remove(p)


def substitute_variables(presentation: PptxPresentation, variables: dict) -> None:
    for slide in presentation.slides:
        expand_variables_on_slide(slide, variables)

    # expand core properties
    presentation.core_properties.author = expand_variables(
        presentation.core_properties.author, variables
    )
    presentation.core_properties.title = expand_variables(
        presentation.core_properties.title, variables
    )
    presentation.core_properties.subject = expand_variables(
        presentation.core_properties.subject, variables
    )


def variable_name(name: str) -> str:
    # we use mustache for templating
    return "{{" + name + "}}"


def expand_variables(text, variables):
    return chevron.render(text, variables)


def expand_variables_on_shape(shape, variables: dict) -> None:
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.runs and paragraph.text:
            for run in paragraph.runs:
                run.text = expand_variables(run.text, variables)

        format_paragraph(paragraph)


def expand_variables_on_slide(
    slide: PptxSlide,
    variables: dict,
    key_value_variables: list = [],
    fit_variables: list = [],
) -> None:
    # substitute shape values first
    for variable, value in variables.items():
        shape = find_shape_by_name(slide, variable_name(variable))
        if shape:
            if is_picture_variable(variable):
                replace_shape_with_picture(slide, shape, value)
            else:
                update_and_format_shape_text(shape, value)
                if variable in key_value_variables:
                    format_key_value_shape(shape)

                if variable in fit_variables:
                    shape.text_frame.fit_text(
                        font_family="Verdana",
                        max_size=10,
                        # line_height=1.1
                    )

    # then expand content
    # expanding content at this point allows us to use templating in values
    for shape in slide.shapes:
        if shape.has_text_frame:
            expand_variables_on_shape(shape, variables)
        elif shape.has_table:
            for r in shape.table.rows:
                for c in r.cells:
                    expand_variables_on_shape(c, variables)


def add_slide(
    deck: PptxPresentation,
    master_layout: SlideLayout,
    variables: dict,
    fit_text_for_variable_names: list = [],
    key_value_variable_names: list = [],
):
    result = deck.slides.add_slide(master_layout)

    expand_variables_on_slide(
        result, variables, key_value_variable_names, fit_text_for_variable_names
    )

    return result
